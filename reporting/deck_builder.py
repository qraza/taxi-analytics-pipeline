"""Programmatic PowerPoint "board pack" builder.

Standalone — no Streamlit imports — so it can be called from the dashboard,
the CLI, or a scheduler alike.
"""

import datetime
import os
from io import BytesIO
from pathlib import Path

import duckdb
import pandas as pd
import plotly.io as pio
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

import cli.llm as llm
from reporting.figures import build_daily_trend_fig, build_hourly_heatmap_fig, build_top_zones_fig

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

ACCENT = RGBColor(0x4C, 0x78, 0xA8)
DARK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x6E, 0x6E, 0x6E)
CARD_FILL = RGBColor(0xF5, 0xF7, 0xFA)
CARD_LINE = RGBColor(0xE0, 0xE4, 0xE9)
FONT = "Calibri"

CHART_LEFT = Inches(0.5)
CHART_TOP = Inches(1.35)
CHART_WIDTH_IN = 12.33
CHART_HEIGHT_IN = 5.65


# ---------------------------------------------------------------------------
# Data access — read-only queries against the marts
# ---------------------------------------------------------------------------

def _load_month_data(conn: duckdb.DuckDBPyConnection, month: str):
    kpis_row = conn.execute(
        """
        SELECT
            sum(total_trips)                                                   as total_trips,
            sum(total_revenue_usd)                                             as total_revenue_usd,
            sum(avg_fare_usd * total_trips) / nullif(sum(total_trips), 0)      as avg_fare_usd,
            sum(avg_duration_minutes * total_trips) / nullif(sum(total_trips), 0) as avg_duration_minutes,
            sum(airport_trip_share * total_trips) / nullif(sum(total_trips), 0)   as airport_trip_share
        FROM main.mart_daily_kpis
        WHERE strftime(trip_date, '%Y-%m') = ?
        """,
        [month],
    ).fetchone()
    kpis = dict(zip(
        ["total_trips", "total_revenue_usd", "avg_fare_usd", "avg_duration_minutes", "airport_trip_share"],
        kpis_row,
    ))

    daily = conn.execute(
        """
        SELECT trip_date, total_trips, total_revenue_usd
        FROM main.mart_daily_kpis
        WHERE strftime(trip_date, '%Y-%m') = ?
        ORDER BY trip_date
        """,
        [month],
    ).df()

    top_zones = conn.execute(
        """
        SELECT
            pickup_zone,
            pickup_borough,
            sum(total_revenue_usd) as total_revenue_usd,
            sum(total_trips)       as total_trips
        FROM main.mart_trip_summary
        WHERE strftime(trip_date, '%Y-%m') = ?
        GROUP BY 1, 2
        ORDER BY total_revenue_usd DESC
        LIMIT 10
        """,
        [month],
    ).df()

    hourly = conn.execute(
        "SELECT * FROM main.mart_hourly_patterns WHERE pickup_month = ?",
        [month],
    ).df()

    return kpis, daily, top_zones, hourly


# ---------------------------------------------------------------------------
# AI commentary — reuses the shared LLM helper, never called unless a key
# is available
# ---------------------------------------------------------------------------

def _generate_insights(
    kpis: dict, top_zones: pd.DataFrame, month_label: str, api_key: str | None
) -> list[str] | None:
    effective_key = api_key or os.environ.get("ANTHROPIC_API_KEY")
    if not effective_key:
        return None

    zones_text = "\n".join(
        f"- {r.pickup_zone} ({r.pickup_borough}): ${r.total_revenue_usd:,.2f} revenue, "
        f"{r.total_trips:,.0f} trips"
        for r in top_zones.itertuples()
    )
    prompt = (
        "You are a data analyst preparing a monthly board report for NYC Yellow Taxi "
        f"operations covering {month_label}.\n\n"
        f"Total trips: {kpis['total_trips']:,.0f}\n"
        f"Total revenue: ${kpis['total_revenue_usd']:,.2f}\n"
        f"Average fare: ${kpis['avg_fare_usd']:,.2f}\n"
        f"Average trip duration: {kpis['avg_duration_minutes']:.1f} minutes\n"
        f"Airport trip share: {kpis['airport_trip_share']:.1%}\n\n"
        f"Top revenue zones:\n{zones_text}\n\n"
        "Write 3-4 crisp bullet points highlighting the most important takeaways for "
        "a board audience. Start each bullet on its own line with '- '. No preamble, "
        "no closing remarks."
    )

    # Temporarily override the shared helper's module-level key so an explicit
    # api_key argument takes precedence over the environment for this call.
    original_key = llm.ANTHROPIC_API_KEY
    llm.ANTHROPIC_API_KEY = effective_key
    try:
        text = llm.call_claude(prompt, max_tokens=300)
    except Exception:
        return None
    finally:
        llm.ANTHROPIC_API_KEY = original_key

    bullets = [
        line.lstrip("-• ").strip()
        for line in text.splitlines()
        if line.strip().startswith(("-", "•"))
    ]
    return bullets or [text.strip()]


# ---------------------------------------------------------------------------
# Slide-building helpers
# ---------------------------------------------------------------------------

def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _style_run(run, size=18, bold=False, color=DARK):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT


def _add_title_bar(slide, title: str):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.7))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    _style_run(run, size=26, bold=True, color=ACCENT)

    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(2.0), Pt(3))
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background()
    rule.shadow.inherit = False


def _add_footer(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    _style_run(run, size=9, color=GRAY)


def _fig_to_png_bytes(fig, width=1600, height=733, scale=2) -> bytes:
    return pio.to_image(fig, format="png", width=width, height=height, scale=scale)


def _add_chart_slide(prs: Presentation, title: str, fig, footer_text: str):
    slide = _blank_slide(prs)
    _add_title_bar(slide, title)
    img_bytes = _fig_to_png_bytes(fig)
    slide.shapes.add_picture(
        BytesIO(img_bytes), CHART_LEFT, CHART_TOP,
        width=Inches(CHART_WIDTH_IN), height=Inches(CHART_HEIGHT_IN),
    )
    _add_footer(slide, footer_text)
    return slide


def _add_kpi_tile(slide, left, top, width, height, value: str, label: str):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_FILL
    card.line.color.rgb = CARD_LINE
    card.line.width = Pt(0.75)
    card.shadow.inherit = False

    tf = card.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.1)

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = value
    _style_run(run1, size=30, bold=True, color=ACCENT)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    _style_run(run2, size=12, color=GRAY)


def _add_insights_section(slide, bullets: list[str], top):
    box = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.3), Inches(2.6))
    tf = box.text_frame
    tf.word_wrap = True

    heading = tf.paragraphs[0]
    run = heading.add_run()
    run.text = "Key Insights"
    _style_run(run, size=16, bold=True, color=ACCENT)

    for bullet in bullets:
        p = tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = f"•  {bullet}"
        _style_run(run, size=13, color=DARK)


def _build_title_slide(prs: Presentation, month_label: str, generated_on: str):
    slide = _blank_slide(prs)

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.15))
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT
    band.line.fill.background()
    band.shadow.inherit = False

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(11.7), Inches(1.3))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = f"NYC Taxi Performance — {month_label}"
    _style_run(run, size=40, bold=True, color=DARK)

    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.9), Inches(11.7), Inches(0.6))
    p2 = subtitle_box.text_frame.paragraphs[0]
    run2 = p2.add_run()
    run2.text = f"Generated {generated_on}  ·  Source: NYC TLC Yellow Taxi trip records"
    _style_run(run2, size=16, color=GRAY)

    return slide


def _build_executive_summary_slide(
    prs: Presentation, kpis: dict, month_label: str,
    insights: list[str] | None, ai_requested: bool,
):
    slide = _blank_slide(prs)
    _add_title_bar(slide, f"Executive Summary — {month_label}")

    tiles = [
        (f"{kpis['total_trips']:,.0f}", "Total Trips"),
        (f"${kpis['total_revenue_usd']:,.0f}", "Total Revenue"),
        (f"${kpis['avg_fare_usd']:,.2f}", "Avg Fare"),
        (f"{kpis['avg_duration_minutes']:.1f} min", "Avg Duration"),
    ]
    tile_width = Inches(2.85)
    tile_height = Inches(1.5)
    gap_in = 0.22
    top = Inches(1.6)
    for i, (value, label) in enumerate(tiles):
        left = Inches(0.5 + i * (2.85 + gap_in))
        _add_kpi_tile(slide, left, top, tile_width, tile_height, value, label)

    airport_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.35), Inches(12.0), Inches(0.45))
    p = airport_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = f"✈  Airport trips: {kpis['airport_trip_share']:.1%} of total volume"
    _style_run(run, size=15, color=DARK)

    if insights:
        _add_insights_section(slide, insights, top=Inches(4.0))
    elif ai_requested:
        note_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(8.0), Inches(0.4))
        p = note_box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "AI commentary unavailable"
        _style_run(run, size=11, color=GRAY)

    _add_footer(slide, "Source: NYC TLC trip records → dbt/DuckDB pipeline")
    return slide


def _build_closing_slide(prs: Presentation):
    slide = _blank_slide(prs)
    _add_title_bar(slide, "Methodology & Sources")

    box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(11.5), Inches(2.0))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Pipeline: TLC trip records → dbt/DuckDB → automated deck"
    _style_run(run, size=18, color=DARK)

    p2 = tf.add_paragraph()
    p2.space_before = Pt(10)
    run2 = p2.add_run()
    run2.text = "Repository: <repo-url-placeholder>"
    _style_run(run2, size=14, color=GRAY)

    return slide


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def build_deck(
    db_path: str,
    month: str = "2024-01",
    include_ai_commentary: bool = False,
    api_key: str | None = None,
) -> BytesIO:
    """Build a monthly board-pack deck from the marts and return it as a BytesIO."""
    if not Path(db_path).exists():
        raise FileNotFoundError(f"DuckDB database not found at {db_path!r}")

    try:
        month_label = datetime.datetime.strptime(month, "%Y-%m").strftime("%B %Y")
    except ValueError as exc:
        raise ValueError(f"month must be in YYYY-MM format, got {month!r}") from exc

    conn = duckdb.connect(db_path, read_only=True)
    try:
        try:
            kpis, daily, top_zones, hourly = _load_month_data(conn, month)
        except duckdb.CatalogException as exc:
            raise RuntimeError(
                "Required marts not found — run `dbt build` before generating a report."
            ) from exc
    finally:
        conn.close()

    if kpis["total_trips"] is None:
        raise ValueError(f"No data found for month {month!r}")

    insights = _generate_insights(kpis, top_zones, month_label, api_key) if include_ai_commentary else None

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    _build_title_slide(prs, month_label, datetime.date.today().strftime("%B %d, %Y"))
    _build_executive_summary_slide(prs, kpis, month_label, insights, include_ai_commentary)
    _add_chart_slide(
        prs, "Daily Trend", build_daily_trend_fig(daily, title=""),
        "Daily trip volume and revenue across the month",
    )
    _add_chart_slide(
        prs, "Top 10 Zones by Revenue", build_top_zones_fig(top_zones, title=""),
        "Pickup zones ranked by total revenue",
    )
    _add_chart_slide(
        prs, "Demand Heatmap", build_hourly_heatmap_fig(hourly, title=""),
        "Trip volume by day of week and pickup hour",
    )
    _build_closing_slide(prs)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
